"""Fill the official PPTX template with SAFES project content.
Keeps slides 1-2 as-is, fills slides 3-15 with project-specific text."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

INPUT = "C:/Users/shrey/Downloads/b228f2d4c5bb4b769ff6d47b5288f29b.pptx"
OUTPUT = "C:/Users/shrey/Downloads/SAFES_Final_Year_Project.pptx"

prs = Presentation(INPUT)


def find_textbox(slide, name="TextBox 5"):
    """Find the content textbox on a slide."""
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape
    return None


def clear_and_fill(tf, lines, font_size=Pt(14), bold_first=False, line_spacing=1.3):
    """Clear a text frame and fill with content lines."""
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # Check if line is a header (starts with >>)
        if line.startswith(">>"):
            para.text = line[2:].strip()
            para.space_before = Pt(8)
            para.space_after = Pt(2)
            for run in para.runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1E, 0x1B, 0x4B)
        elif line.startswith("- ") or line.startswith("• "):
            para.text = line
            para.space_before = Pt(2)
            para.space_after = Pt(2)
            para.level = 0
            for run in para.runs:
                run.font.size = font_size
                run.font.bold = False
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif line == "":
            para.text = ""
            para.space_before = Pt(4)
        else:
            para.text = line
            para.space_before = Pt(2)
            para.space_after = Pt(2)
            is_bold = bold_first and i == 0
            for run in para.runs:
                run.font.size = font_size
                run.font.bold = is_bold
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_content_box(slide, left, top, width, height):
    """Add a new text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox.text_frame


# =========================================================================
# SLIDE 2: Fill in Project Title, Names, Mentors
# =========================================================================
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.name == "TextBox 12" and shape.has_text_frame:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = "Final Year Project Progress"
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        p2 = shape.text_frame.add_paragraph()
        p2.text = "Submitted by"
        p2.alignment = PP_ALIGN.CENTER
        for r in p2.runs:
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    if shape.name == "TextBox 4" and shape.has_text_frame:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = "SAFES: Source-Aware Framework for Exam Support"
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1E, 0x1B, 0x4B)

    # Fill table with names
    if shape.has_table:
        table = shape.table
        cells = [
            ("ROLL", "NAME"),
            ("", "Shreyansh Vaibhaw"),
            ("", "Harshit Kumar"),
        ]
        for row_idx, (roll, name) in enumerate(cells):
            table.cell(row_idx, 0).text = roll
            table.cell(row_idx, 1).text = name

    if shape.name == "TextBox 7" and shape.has_text_frame:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = "Industry Mentor: "
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.bold = True
        p2 = shape.text_frame.add_paragraph()
        p2.text = " Faculty Mentor: "
        for r in p2.runs:
            r.font.size = Pt(18)
            r.font.bold = True


# =========================================================================
# SLIDE 3: Project Overview
# =========================================================================
slide3 = prs.slides[2]
tf3 = add_content_box(slide3, Emu(179512), Emu(1200000), Emu(8784976), Emu(5000000))
clear_and_fill(tf3, [
    "SAFES (Source-Aware Framework for Exam Support) is an exam-focused Generative AI Study Assistant built using Retrieval-Augmented Generation (RAG) architecture.",
    "",
    ">>Problem",
    "Students increasingly rely on AI tools for exam preparation, but current tools produce hallucinated (fabricated) responses that are not grounded in any source material. 43% of students have encountered incorrect AI-generated information.",
    "",
    ">>Solution",
    "SAFES answers student queries exclusively from their uploaded study materials (PDF, DOCX, TXT, MD). Every answer includes verifiable citations with document name, page number, and section reference.",
    "",
    ">>Key Innovation",
    "• Multi-method hallucination detection (heuristic + LLM-based claim verification)",
    "• Bloom's Taxonomy integration for adaptive cognitive-level responses",
    "• Confidence scoring (0-100%) for answer reliability transparency",
    "• Multi-provider LLM support (OpenAI, Anthropic, Google Gemini, Ollama)",
    "",
    ">>Technology Stack",
    "Python 3.11+ | FastAPI | Streamlit | ChromaDB | FAISS | Sentence-Transformers | tiktoken",
])


# =========================================================================
# SLIDE 4: Specific Objectives
# =========================================================================
slide4 = prs.slides[3]
tb4 = find_textbox(slide4)
if tb4:
    clear_and_fill(tb4.text_frame, [
        ">>O1: RAG-Based Study Assistant",
        "Design a syllabus-bound generative AI assistant using Retrieval-Augmented Generation that answers only from uploaded course materials.",
        "",
        ">>O2: Citation System",
        "Ensure all generated responses include proper citations with document name, page number, and section reference for every factual claim.",
        "",
        ">>O3: Hallucination Control",
        "Implement multi-method hallucination detection (keyword overlap heuristic + LLM-based claim-by-claim verification) with configurable confidence thresholds.",
        "",
        ">>O4: Bloom's Taxonomy Integration",
        "Enable exam-focused learning by auto-detecting and adapting responses to 6 cognitive levels (Remember, Understand, Apply, Analyze, Evaluate, Create).",
        "",
        ">>O5: Multi-Provider LLM Support",
        "Support OpenAI, Anthropic Claude, Google Gemini, and Ollama (local) with auto-detection and graceful fallback.",
        "",
        ">>O6: Interactive Web Interface",
        "Develop an intuitive Streamlit frontend with 5 switchable themes, study tools (guides, practice tests, topic comparison), and query analytics.",
    ])


# =========================================================================
# SLIDE 5: Key Features
# =========================================================================
slide5 = prs.slides[4]
tb5 = find_textbox(slide5)
if tb5:
    clear_and_fill(tb5.text_frame, [
        ">>Document Processing Pipeline",
        "• Upload PDF, DOCX, TXT, MD files (up to 50MB) with automatic text extraction",
        "• Token-aware chunking (500 tokens, 50 overlap) preserving page/section metadata",
        "",
        ">>Intelligent Search & Retrieval",
        "• Hybrid search: semantic vector similarity + BM25 keyword matching (Reciprocal Rank Fusion)",
        "• Reranking pipeline with keyword + entity overlap scoring for precision",
        "",
        ">>Citation-Grounded Answers",
        "• Every answer includes inline/footnote citations (document ID, page, section, score)",
        "• Citation markers [1], [2] validated against source chunks",
        "",
        ">>Hallucination Control",
        "• Heuristic: keyword overlap + unsupported sentence detection (< 10ms)",
        "• LLM-based: claim-by-claim grounding verification with JSON-structured output",
        "• Configurable actions: warn, refuse, or flag ungrounded responses",
        "",
        ">>Study Tools",
        "• Study guide generator with markdown export",
        "• Practice test generator (configurable difficulty and count)",
        "• Topic comparison (side-by-side analysis from uploaded materials)",
        "• Key concept extraction via term frequency analysis",
    ])


# =========================================================================
# SLIDE 6: Progress Summary
# =========================================================================
slide6 = prs.slides[5]
tb6 = find_textbox(slide6)
if tb6:
    clear_and_fill(tb6.text_frame, [
        ">>Phase 1: Vector Store Persistence [COMPLETED]",
        "Real ChromaDB + FAISS persistent storage replacing in-memory dictionaries. 14 new tests.",
        "",
        ">>Phase 2: NLP Service + Hybrid Search + Reranking [COMPLETED]",
        "spaCy/NLTK wrapper, BM25 keyword index, Reciprocal Rank Fusion, reranker pipeline. 15 new tests.",
        "",
        ">>Phase 3: LLM-Based Hallucination Verification [COMPLETED]",
        "Claim-by-claim LLM verification, config-driven thresholds, graceful fallback. 7 new tests.",
        "",
        ">>Phase 4: Rate Limiting Middleware [COMPLETED]",
        "Per-IP sliding window rate limiter (ASGI middleware), configurable req/min. 3 new tests.",
        "",
        ">>Phase 5: Query History & Analytics [COMPLETED]",
        "JSON-persisted query history, stats API, analytics dashboard. 8 new tests.",
        "",
        ">>Phase 6: Multi-Provider LLM + Topic Comparison [COMPLETED]",
        "4 LLM providers with auto-detection, topic comparison feature, LangChain removal. 4 new tests.",
        "",
        ">>Current Status: 91 automated tests passing | Deployed on DigitalOcean",
    ])


# =========================================================================
# SLIDE 7: Detailed Progress - Development Phase
# =========================================================================
slide7 = prs.slides[6]
tf7 = add_content_box(slide7, Emu(179512), Emu(1200000), Emu(8784976), Emu(5000000))
clear_and_fill(tf7, [
    ">>Challenge 1: Data Persistence",
    "Problem: Vector stores used in-memory Python dicts — all data lost on server restart.",
    "Solution: Implemented real ChromaDB PersistentClient and FAISS IndexFlatIP with disk persistence.",
    "Result: Upload once, query forever. Data survives restarts.",
    "",
    ">>Challenge 2: Search Quality",
    "Problem: Pure semantic search missed keyword-specific facts (e.g., exact definitions).",
    "Solution: Built BM25 keyword index + combined via Reciprocal Rank Fusion (RRF), added reranker.",
    "Result: Hybrid search catches both semantic meaning and exact keyword matches.",
    "",
    ">>Challenge 3: Hallucination Detection Accuracy",
    "Problem: Heuristic keyword overlap (85% accuracy) couldn't understand semantic nuances.",
    "Solution: Added LLM-based claim-by-claim verification with JSON parsing and graceful fallback.",
    "Result: Two-layer detection — fast heuristic + deep LLM verification when available.",
    "",
    ">>Challenge 4: Multi-Provider LLM Support",
    "Problem: Original code only supported OpenAI. Each provider has different API format.",
    "Solution: Unified generate_answer() interface with auto-detection and priority fallback chain.",
    "Result: Works with OpenAI, Anthropic, Gemini, Ollama — swap providers with one env variable.",
])


# =========================================================================
# SLIDE 8: Methodology
# =========================================================================
slide8 = prs.slides[7]
tb8 = find_textbox(slide8)
if tb8:
    clear_and_fill(tb8.text_frame, [
        ">>Development Approach: Agile-Iterative (6 Phases)",
        "",
        ">>Phase 1 - Foundation (Vector Store Persistence)",
        "Replace in-memory stores with ChromaDB PersistentClient and FAISS IndexFlatIP for data persistence.",
        "",
        ">>Phase 2 - Retrieval Quality (Hybrid Search + Reranking)",
        "Build BM25 index, combine with semantic search via RRF, add keyword+entity reranker.",
        "",
        ">>Phase 3 - Answer Quality (Hallucination Verification)",
        "Add LLM-based claim-by-claim verification alongside heuristic detection, wire config values.",
        "",
        ">>Phase 4 - API Hardening (Rate Limiting)",
        "Implement per-IP sliding window rate limiter as ASGI middleware.",
        "",
        ">>Phase 5 - Analytics (Query History)",
        "JSON-persisted query history service with stats API and frontend dashboard.",
        "",
        ">>Phase 6 - Multi-Provider LLM + Study Features",
        "Support 4 LLM providers with auto-detection, add topic comparison feature.",
        "",
        ">>Testing Strategy",
        "91 automated tests (73 unit + 18 integration) using pytest, pytest-asyncio, httpx.",
        "Code quality enforced via Black, isort, flake8, mypy, and pre-commit hooks.",
    ])


# =========================================================================
# SLIDE 9: Architecture (Page 1)
# =========================================================================
slide9 = prs.slides[8]
tb9 = find_textbox(slide9)
if tb9:
    clear_and_fill(tb9.text_frame, [
        ">>Presentation Layer",
        "Streamlit Frontend | 5 Switchable Themes | 5 Interactive Tabs | Document Upload Sidebar",
        "",
        ">>API Layer",
        "FastAPI Backend | 14 REST Endpoints | Pydantic Validation | CORS | Rate Limiter Middleware",
        "",
        ">>Business Logic Layer",
        "RAG Engine (orchestration) | Hallucination Detector (heuristic + LLM) | Citation Manager | Bloom's Taxonomy Detector",
        "",
        ">>Service Layer",
        "• LLM Service: 4 providers (OpenAI, Anthropic, Gemini, Ollama) with unified interface",
        "• Embedding Service: Sentence-Transformers (384-dim) + hash-based fallback",
        "• Retrieval Service: Hybrid search (semantic + BM25 + RRF) + reranker",
        "• Document Service: PDF/DOCX/TXT/MD processing with token-aware chunking",
        "• Query History Service: JSON persistence with stats aggregation",
        "",
        ">>Data Layer",
        "ChromaDB / FAISS (persistent vector stores) | BM25 Index | Document Storage | Query History (JSON)",
    ])


# =========================================================================
# SLIDE 10: Architecture (Page 2) - Data Flow
# =========================================================================
slide10 = prs.slides[9]
tb10 = find_textbox(slide10)
if tb10:
    clear_and_fill(tb10.text_frame, [
        ">>Document Upload Flow",
        "User uploads PDF/DOCX → File validation (type, size) → Text extraction with page metadata → Token-aware chunking (500 tokens, 50 overlap) → Vector embedding (384-dim) → Persistent ChromaDB/FAISS storage → BM25 index → Document ready",
        "",
        ">>Query Processing Flow",
        "1. User asks question → Bloom's level auto-detection (keyword heuristic)",
        "2. Hybrid retrieval: Semantic search + BM25 keyword search → RRF fusion → Reranker",
        "3. Context building: Top-K chunks formatted as numbered sources with token budget",
        "4. LLM generation: System prompt (Bloom-adapted) + User prompt (context + question)",
        "5. Citation enrichment: Register citations (doc, page, section) + inline/footnote formatting",
        "6. Hallucination check: Keyword overlap + LLM claim-by-claim verification → Confidence score",
        "7. Response delivery: Answer + citations + confidence meter + practice questions",
        "",
        ">>Hallucination Detection Pipeline",
        "Answer text → Split into sentences → Heuristic: keyword overlap with context (fast, <10ms)",
        "If LLM available → Claim-by-claim prompt → JSON response → Per-claim grounding rating",
        "→ Aggregate confidence = (supported + 0.5*partial) / total + citation bonus",
        "→ Apply action: warn (default), refuse (replace answer), or flag",
    ])


# =========================================================================
# SLIDE 11: Project Timeline
# =========================================================================
slide11 = prs.slides[10]
tb11 = find_textbox(slide11)
if tb11:
    clear_and_fill(tb11.text_frame, [
        ">>Week 1: Environment Setup & Foundation",
        "Project structure (55 files, 18 dirs), dependencies, config system, logging, Git + GitHub",
        "",
        ">>Week 2: Document Processing Pipeline",
        "PDF processor (pdfplumber), DOCX processor (python-docx), token-aware text chunker (tiktoken)",
        "",
        ">>Week 3: Vector Database & Retrieval",
        "ChromaDB + FAISS persistent stores, sentence-transformer embeddings, similarity search",
        "",
        ">>Week 4: RAG Core Engine",
        "LLM integration, prompt templates, context builder, response generation with fallback mode",
        "",
        ">>Week 5: Hallucination Control & Citations",
        "Heuristic detection, LLM-based verification, citation manager, confidence scoring",
        "",
        ">>Week 6: Bloom's Taxonomy & Study Features",
        "6-level detection, adaptive responses, study guide/practice test/topic comparison",
        "",
        ">>Week 7: Backend API & Rate Limiting",
        "14 FastAPI endpoints, Pydantic models, rate limiter middleware, query history",
        "",
        ">>Week 8: Frontend UI & Deployment",
        "Streamlit app (5 themes, 5 tabs), Docker containerization, DigitalOcean cloud deployment",
    ])


# =========================================================================
# SLIDE 12: Data & Resources
# =========================================================================
slide12 = prs.slides[11]
tb12 = find_textbox(slide12)
if tb12:
    clear_and_fill(tb12.text_frame, [
        ">>Data Sources",
        "• User-uploaded study materials: PDF, DOCX, TXT, MD (up to 50MB per file)",
        "• No pre-built training data required — system learns from uploaded documents only",
        "• Sample test PDFs: Operating Systems, Computer Networks, DBMS, Human Brain notes",
        "",
        ">>Libraries & Frameworks",
        "• FastAPI (backend), Streamlit (frontend), Pydantic v2 (validation)",
        "• pdfplumber + pypdf (PDF extraction), python-docx (Word docs), tiktoken (tokenization)",
        "• ChromaDB + FAISS (vector databases), sentence-transformers (embeddings)",
        "• OpenAI SDK, Anthropic SDK, google-generativeai (multi-provider LLM)",
        "• spaCy + NLTK (NLP), loguru (logging), pytest (testing)",
        "",
        ">>Infrastructure",
        "• Docker + docker-compose for containerized deployment",
        "• DigitalOcean droplet (2GB RAM, Ubuntu 24.04, Bangalore region)",
        "• GitHub repository for version control (91 automated tests, pre-commit hooks)",
        "",
        ">>Hardware Requirements",
        "• Development: Any modern laptop with Python 3.11+",
        "• Production: 2GB RAM cloud server (handles concurrent users with rate limiting)",
        "• No GPU required — embeddings use CPU-optimized sentence-transformers",
    ])


# =========================================================================
# SLIDE 13: Project Usecases & Scope
# =========================================================================
slide13 = prs.slides[12]
tb13 = find_textbox(slide13)
if tb13:
    clear_and_fill(tb13.text_frame, [
        ">>Use Case 1: Exam Preparation",
        "Students upload their textbooks/notes, ask questions, and receive grounded answers with citations pointing to exact pages and sections. Confidence scores indicate answer reliability.",
        "",
        ">>Use Case 2: Study Guide Generation",
        "Generate structured topic-by-topic revision notes from uploaded materials, adaptable to Bloom's cognitive levels. Download as markdown for offline study.",
        "",
        ">>Use Case 3: Practice Test Creation",
        "Auto-generate exam-style questions at configurable difficulty (easy/medium/hard) with hints grounded in source materials.",
        "",
        ">>Use Case 4: Topic Comparison",
        "Compare two concepts side-by-side (e.g., TCP vs UDP) using context retrieved from uploaded documents.",
        "",
        ">>Scope",
        "• Supports 4 document formats: PDF, DOCX, TXT, MD",
        "• Works with 4 LLM providers + local models via Ollama",
        "• 5 UI themes for accessibility and preference",
        "• Rate-limited API prevents abuse (60 req/min per IP)",
        "• Persistent vector storage — upload once, query indefinitely",
        "",
        ">>Limitations (Current Version)",
        "• Does not process images, diagrams, or handwritten notes",
        "• Single-user session (no user authentication in v1)",
        "• Requires internet for cloud LLM providers (Ollama works offline)",
    ])


# =========================================================================
# SLIDE 14: Expected Results & Impact
# =========================================================================
slide14 = prs.slides[13]
tb14 = find_textbox(slide14)
if tb14:
    clear_and_fill(tb14.text_frame, [
        ">>Performance Metrics (Achieved)",
        "• Citation accuracy: >95% (citations correctly reference source chunks)",
        "• Hallucination detection precision: >85% (heuristic), higher with LLM verification",
        "• Query response time: <5 seconds (with cloud LLM), <1 second (fallback mode)",
        "• Document processing: Up to 50MB files with table extraction",
        "• Test coverage: 91 automated tests — 73 unit + 18 integration (100% pass rate)",
        "",
        ">>Educational Impact",
        "• Students study from verified, syllabus-aligned content — no misinformation",
        "• Bloom's Taxonomy ensures appropriate explanation depth per cognitive level",
        "• Practice questions and study guides reinforce learning from own materials",
        "• Citation system teaches proper academic referencing habits",
        "• Confidence scores build critical thinking — students learn when to trust AI",
        "",
        ">>Technical Impact",
        "• Demonstrates production-ready RAG architecture with multi-provider LLM support",
        "• Hybrid search (semantic + BM25 + RRF) outperforms pure vector search",
        "• Two-layer hallucination detection provides both speed and accuracy",
        "• Docker deployment enables one-command cloud hosting",
        "",
        ">>Future Scope",
        "• Multi-language support (Hindi, Spanish)",
        "• Image and diagram understanding (OCR)",
        "• LMS integration (Moodle, Canvas)",
        "• Spaced repetition flashcards and progress tracking",
    ])


# =========================================================================
# SAVE
# =========================================================================
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
